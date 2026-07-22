"""Generate a 4-quadrant technology commercialization slide for the ORNL IP filing."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Brand colors ---
ORNL_GREEN = RGBColor(0x00, 0x66, 0x2C)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
QUAD_COLORS = [
    RGBColor(0x0B, 0x4F, 0x6C),  # teal
    RGBColor(0x5A, 0x3E, 0x85),  # purple
    RGBColor(0x9C, 0x4A, 0x1E),  # rust
    RGBColor(0x1E, 0x6B, 0x3A),  # green
]

# --- Content: 4 quadrants for the IP filing ---
TITLE = "LLM & RAG-Powered Autonomous Software Development"
SUBTITLE = ("for Creating Scientific Visual Dashboards  |  ORNL Invention Disclosure "
            "202405756  |  Disclosed 29 Aug 2024  |  Status: To Be Filed")

quadrants = [
    {
        "title": "Opportunity / Technology Description",
        "bullets": [
            ("Problem addressed", "Building scientific dashboards is slow, code-heavy, and "
             "requires scarce full-stack + domain expertise."),
            ("Technical description", "An autonomous agent that couples Large Language Models "
             "(LLMs) with Retrieval-Augmented Generation (RAG) to plan, generate, test, and "
             "iteratively refine end-to-end dashboard software from natural-language intent."),
            ("Key components", "RAG knowledge base over datasets, APIs & design patterns; "
             "LLM code-generation pipeline; automated validation/self-correction loop; "
             "visual dashboard renderer."),
            ("Development stage", "Working prototype demonstrated on scientific/visualization "
             "workloads."),
            ("Fundamental advantage", "Grounded (RAG) generation reduces hallucination and "
             "produces deployable, data-accurate dashboards with minimal human coding."),
        ],
    },
    {
        "title": "Market Opportunity & Expected Impact",
        "bullets": [
            ("Who has the problem", "National labs, research institutions, and enterprises "
             "needing rapid, data-driven visual dashboards."),
            ("Current solutions & gap", "BI tools (Tableau, Power BI) and hand-coded dashboards "
             "are rigid, manual, and disconnected from live scientific data pipelines."),
            ("Impact", "Cuts dashboard build time from weeks to hours; democratizes access for "
             "non-programmer domain scientists."),
            ("Applications", "Experiment monitoring, simulation output visualization, "
             "operational/facility dashboards, reporting for sponsors."),
            ("Licensing pathway", "Non-exclusive software license or SaaS; consider "
             "field-limited exclusivity for vertical integrations, coupled to CRADA/SPP pilots."),
        ],
    },
    {
        "title": "Go-to-Market & Technical Development Plan",
        "bullets": [
            ("Months 1-3 (discovery)", "Publish technology listing; identify pilot partners "
             "across labs, universities, and analytics vendors; run NDA-protected demo webinar."),
            ("Months 2-11 (validation)", "Harden RAG grounding on partner datasets; benchmark "
             "generation accuracy & self-correction; validate security, provenance, and "
             "reproducibility of generated code."),
            ("Deliverable", "Partner-ready package: reproducible pipeline, evaluation metrics, "
             "and reference dashboard deployments."),
            ("Key risks", "LLM hallucination/accuracy; data governance & IP of generated code; "
             "integration variance across data stacks; model/version drift."),
        ],
    },
    {
        "title": "Value Proposition & Deployment Plan",
        "bullets": [
            ("Alternatives", "Manual coding, static BI dashboards, generic code copilots without "
             "data grounding."),
            ("Why better", "RAG grounding yields data-accurate, deployable dashboards; autonomous "
             "loop reduces developer effort; adapts to new datasets without rebuild."),
            ("Integration challenges", "Medium: data connectors, access control, on-prem vs "
             "cloud LLM hosting, and validation of generated artifacts."),
            ("Partners / licensees", "Analytics & visualization software vendors, cloud/AI "
             "platform providers, research computing organizations."),
            ("End users", "Domain scientists, data analysts, program managers, and facility "
             "operators."),
        ],
    },
]

# --- Build slide ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

SW, SH = prs.slide_width, prs.slide_height

# Header band
hdr = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.05))
hdr.fill.solid(); hdr.fill.fore_color.rgb = ORNL_GREEN
hdr.line.fill.background()
tf = hdr.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.05)
p = tf.paragraphs[0]
r = p.add_run(); r.text = TITLE
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = SUBTITLE
r2.font.size = Pt(10.5); r2.font.color.rgb = RGBColor(0xE8, 0xF0, 0xE8)

# Quadrant geometry
top = Inches(1.2)
gap = Inches(0.12)
qw = (SW - gap) / 2
qh = (SH - top - gap - Inches(0.05)) / 2
positions = [
    (0, top),
    (qw + gap, top),
    (0, top + qh + gap),
    (qw + gap, top + qh + gap),
]

for (x, y), quad, color in zip(positions, quadrants, QUAD_COLORS):
    box = slide.shapes.add_shape(1, int(x), int(y), int(qw), int(qh))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    box.shadow.inherit = False

    # Title bar
    bar = slide.shapes.add_shape(1, int(x), int(y), int(qw), Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    btf = bar.text_frame; btf.word_wrap = True
    btf.margin_left = Inches(0.15); btf.margin_top = Inches(0.02); btf.margin_bottom = Inches(0.02)
    bp = btf.paragraphs[0]
    br = bp.add_run(); br.text = quad["title"]
    br.font.size = Pt(13); br.font.bold = True; br.font.color.rgb = WHITE

    # Body text
    tb = slide.shapes.add_textbox(int(x) + Inches(0.15), int(y) + Inches(0.5),
                                  int(qw) - Inches(0.3), int(qh) - Inches(0.55))
    btf = tb.text_frame; btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (label, text) in enumerate(quad["bullets"]):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.space_after = Pt(4)
        lr = para.add_run(); lr.text = f"{label}: "
        lr.font.size = Pt(10); lr.font.bold = True; lr.font.color.rgb = color
        tr = para.add_run(); tr.text = text
        tr.font.size = Pt(10); tr.font.color.rgb = DARK

out = "/home/jose/nsrd_ornl/nsrd_ui/ttm/IP_202405756_4quadrant.pptx"
prs.save(out)
print("Saved:", out)
